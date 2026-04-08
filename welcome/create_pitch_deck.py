"""
AI Badge 90-Second Pitch Deck Generator
Creates a compelling 8-slide pitch deck using the AI Badge template.
Victor records his voice over these slides for a 90-second intro video.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
import os

# Brand colors
NAVY = RGBColor(0x00, 0x00, 0x36)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xC8, 0xC8, 0xE0)
DARK_TEXT = RGBColor(0x0A, 0x0E, 0x27)
SOFT_GOLD = RGBColor(0xB8, 0x96, 0x0E)

# Slide dimensions (16:9)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), '..', 'tutorials', 'hello-world-2', 'hello-world-2.pptx')
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), 'ai-badge-pitch.pptx')
BADGE_IMG = os.path.join(os.path.dirname(__file__), '..', 'assets', 'img', 'ai-badge-credential.png')


def add_dark_bg(slide):
    """Add navy gradient background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_light_bg(slide):
    """Add light background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)


def add_text_box(slide, left, top, width, height, text, font_name='Arial', font_size=24,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_name='Arial', font_size=20,
                       color=WHITE, alignment=PP_ALIGN.LEFT, line_spacing=1.4):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.alignment = alignment
        p.line_spacing = Pt(font_size * line_spacing)
        p.space_after = Pt(font_size * 0.4)
    return txBox


def add_gold_bar(slide, top, width=Inches(2)):
    """Add a thin gold accent bar."""
    left = (SLIDE_W - width) // 2
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_stat_block(slide, left, top, number, label):
    """Add a stat number + label vertically."""
    add_text_box(slide, left, top, Inches(2.8), Inches(0.8),
                 number, 'Arial Black', 44, True, GOLD, PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.75), Inches(2.8), Inches(0.6),
                 label, 'Arial', 14, False, LIGHT_GREY, PP_ALIGN.CENTER)


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # We'll use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # ========================================
    # SLIDE 1: HOOK (10s)
    # "AI Literacy Is No Longer Optional"
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide)

    add_text_box(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.5),
                 'AI BADGE', 'Arial Black', 18, True, GOLD, PP_ALIGN.LEFT)

    add_gold_bar(slide, Inches(1.7), Inches(1.5))

    add_text_box(slide, Inches(1.2), Inches(2.0), Inches(9.5), Inches(2.5),
                 'AI Literacy Is No Longer\nOptional.',
                 'Arial Black', 52, True, WHITE, PP_ALIGN.LEFT, 1.15)

    add_text_box(slide, Inches(1.2), Inches(4.5), Inches(9), Inches(1.0),
                 'It is a career requirement.',
                 'Arial', 28, False, GOLD, PP_ALIGN.LEFT)

    # ========================================
    # SLIDE 2: AGITATE (10s)
    # Stats that prove the urgency
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8),
                 'The gap is widening. Every week.',
                 'Arial', 22, False, LIGHT_GREY, PP_ALIGN.CENTER)

    add_gold_bar(slide, Inches(1.5), Inches(3))

    # Three stats
    add_stat_block(slide, Inches(0.5), Inches(2.0), '92%', 'of executives say AI skills\nwill be required within 18 months')
    add_stat_block(slide, Inches(4.35), Inches(2.0), '68%', 'of hiring managers now filter\ncandidates by AI capability')
    add_stat_block(slide, Inches(8.2), Inches(2.0), '3x', 'faster promotion for professionals\nwith verified AI competency')

    add_text_box(slide, Inches(1.5), Inches(5.2), Inches(9), Inches(0.8),
                 'YouTube tutorials won\'t get you there.\nYou need a structured path, a credential, and proof.',
                 'Arial', 18, False, LIGHT_GREY, PP_ALIGN.CENTER, 1.4)

    # ========================================
    # SLIDE 3: SOLUTION (10s)
    # What AI Badge IS
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
                 'AI BADGE', 'Arial Black', 18, True, GOLD, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10.4), Inches(1.5),
                 'Go From AI Curious\nto AI Fluent in Six Weeks.',
                 'Arial Black', 44, True, WHITE, PP_ALIGN.CENTER, 1.15)

    add_gold_bar(slide, Inches(3.1), Inches(2.5))

    add_text_box(slide, Inches(1.5), Inches(3.5), Inches(9), Inches(1.5),
                 'Personalised, one-to-one AI coaching.\nWe measure where you are, guide you forward,\nand prove how far you\'ve come.',
                 'Arial', 22, False, LIGHT_GREY, PP_ALIGN.CENTER, 1.5)

    add_text_box(slide, Inches(2), Inches(5.5), Inches(8), Inches(0.5),
                 'Six weeks.  Six dimensions.  Real results.',
                 'Arial', 20, True, GOLD, PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 4: HOW IT WORKS (15s)
    # The 4-step journey
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide)

    add_text_box(slide, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8),
                 'Your Journey', 'Arial Black', 36, True, WHITE, PP_ALIGN.CENTER)

    add_gold_bar(slide, Inches(1.4), Inches(2.5))

    steps = [
        ('01', 'ASSESS', 'Map your capability\nacross six dimensions'),
        ('02', 'CALIBRATE', 'Get a personalised\ndevelopment pathway'),
        ('03', 'COACH', 'Weekly 1:1 check-ins\nwith expert guidance'),
        ('04', 'PROVE', 'Measurable before-\nand-after progression'),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.6 + i * 2.9)
        y = Inches(2.0)

        add_text_box(slide, x, y, Inches(2.5), Inches(0.6),
                     num, 'Arial Black', 36, True, GOLD, PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.7), Inches(2.5), Inches(0.5),
                     title, 'Arial Black', 18, True, WHITE, PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(1.3), Inches(2.5), Inches(0.8),
                     desc, 'Arial', 14, False, LIGHT_GREY, PP_ALIGN.CENTER, 1.5)

    add_text_box(slide, Inches(1.5), Inches(5.3), Inches(9), Inches(0.6),
                 '1-2 hours per week  |  Designed for working professionals',
                 'Arial', 16, False, LIGHT_GREY, PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 5: WHAT YOU GET (10s)
    # Credential + skills
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.5), Inches(0.8),
                 'What You Earn', 'Arial Black', 36, True, WHITE, PP_ALIGN.LEFT)

    # Left side: benefits list
    benefits = [
        '\u2713  Full AI Competency Assessment (before + after)',
        '\u2713  Personalised development pathway',
        '\u2713  6 weekly one-to-one coaching sessions',
        '\u2713  Guided exercises calibrated to your level',
        '\u2713  Tutorial library access',
        '\u2713  Progress tracking dashboard',
        '\u2713  Verified digital credential (LinkedIn-ready)',
    ]

    add_multiline_text(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.0),
                       benefits, 'Arial', 16, LIGHT_GREY, PP_ALIGN.LEFT, 1.6)

    # Right side: badge image
    if os.path.exists(BADGE_IMG):
        slide.shapes.add_picture(BADGE_IMG, Inches(7.0), Inches(1.2), Inches(4.5), Inches(4.5))

    add_text_box(slide, Inches(6.8), Inches(5.8), Inches(5), Inches(0.5),
                 'Your proof of AI fluency.', 'Arial', 16, True, GOLD, PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 6: SOCIAL PROOF (10s)
    # Instructor + credibility
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide)

    add_text_box(slide, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8),
                 'Your Coach', 'Arial Black', 36, True, WHITE, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(9), Inches(0.5),
                 'VICTOR DEL ROSAL', 'Arial Black', 24, True, GOLD, PP_ALIGN.CENTER)

    credentials = [
        'Award-winning AI educator & author',
        'Chief AI Officer, fiveinnolabs',
        'Adjunct Professor of AI (NCI, Master\'s level)',
        'Trained Fortune 500 leaders on GenAI',
        'Harvard, Oxford, UCD Smurfit, Tec de Monterrey',
        '170+ GitHub repositories (builds what he teaches)',
        'Teaching Award, National Forum for T&L Enhancement',
    ]

    add_multiline_text(slide, Inches(2), Inches(2.3), Inches(8), Inches(3.5),
                       credentials, 'Arial', 17, LIGHT_GREY, PP_ALIGN.CENTER, 1.5)

    add_text_box(slide, Inches(1), Inches(5.7), Inches(10), Inches(0.6),
                 '"Whether you are new to AI or at an intermediate level,\nyou will feel guided, challenged, and supported."',
                 'Arial', 16, False, GOLD, PP_ALIGN.CENTER, 1.4)

    # ========================================
    # SLIDE 7: GUARANTEE (10s)
    # Risk reversal
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide)

    add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(1.0),
                 '100% Guarantee', 'Arial Black', 48, True, GOLD, PP_ALIGN.CENTER)

    add_gold_bar(slide, Inches(2.5), Inches(3))

    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(9), Inches(1.8),
                 'If you have completed all sessions and all exercises,\nand still do not find a single practical use case,\nwe will refund your full tuition.',
                 'Arial', 24, False, WHITE, PP_ALIGN.CENTER, 1.6)

    add_text_box(slide, Inches(2), Inches(5.0), Inches(8), Inches(0.6),
                 'No forms. No friction. Just email me for a full refund.',
                 'Arial', 22, True, GOLD, PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 8: CTA (15s)
    # Pick a slot, start now
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.8), Inches(10), Inches(0.5),
                 'AI BADGE', 'Arial Black', 20, True, GOLD, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(11), Inches(1.5),
                 'Your Starting Point\nIs Right Here.',
                 'Arial Black', 48, True, WHITE, PP_ALIGN.CENTER, 1.15)

    add_gold_bar(slide, Inches(3.6), Inches(2.5))

    # Pricing
    add_text_box(slide, Inches(2), Inches(4.0), Inches(8), Inches(0.6),
                 'From \u20ac90/week  |  6 weeks  |  One-to-one coaching',
                 'Arial', 20, False, LIGHT_GREY, PP_ALIGN.CENTER)

    # CTA button shape
    btn_w = Inches(4)
    btn_h = Inches(0.8)
    btn_left = (SLIDE_W - btn_w) // 2
    btn_top = Inches(5.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, btn_left, btn_top, btn_w, btn_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'Pick a Slot'
    p.font.name = 'Arial Black'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, Inches(2), Inches(6.0), Inches(8), Inches(0.4),
                 'aibadge.fiveinnolabs.com',
                 'Arial', 16, False, GOLD, PP_ALIGN.CENTER)

    # Save
    prs.save(OUTPUT_PATH)
    print(f'Pitch deck saved to: {OUTPUT_PATH}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    build_deck()
