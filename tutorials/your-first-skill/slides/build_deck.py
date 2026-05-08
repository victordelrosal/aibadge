#!/usr/bin/env python3
"""
Build the Your First Skill slide deck (Lesson 9).

Modelled on tutorials/your-ai-team/slides/build_deck.py. 16:9, locked
Indigo Royal palette. White `main` body slides. Hero images preserve native
aspect. Per feedback_slide_design.md: blue cover, hero-dominant, no
rhetorical close, soften absolutes, minimal card decoration.

Run: python3 build_deck.py
Output: ../your-first-skill.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
OUT = HERE.parent / "your-first-skill.pptx"
COVER_TEMPLATE = HERE.parent.parent / "cover" / "cover.pptx"

# --- Locked Indigo Royal palette (light mode) ---
INDIGO = RGBColor(0x1e, 0x1d, 0x4c)
INDIGO_SOFT = RGBColor(0x2b, 0x2e, 0x45)
INDIGO_DIM = RGBColor(0x5b, 0x62, 0x78)
ROYAL = RGBColor(0x00, 0x33, 0x99)
ROYAL_SOFT = RGBColor(0x6f, 0x82, 0xd4)
PANEL_LIGHT = RGBColor(0xf6, 0xf7, 0xfb)
HAIRLINE = RGBColor(0xd8, 0xdc, 0xea)
WHITE = RGBColor(0xff, 0xff, 0xff)
GOLD = RGBColor(0xc1, 0x8a, 0x1c)
GREEN_OK = RGBColor(0x2f, 0x6b, 0x3d)
RED_BAD = RGBColor(0x8a, 0x2a, 0x3f)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _layout_by_name(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[-1]


def add_main_slide(prs, title_text):
    layout = _layout_by_name(prs, "main")
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name = "Helvetica Neue"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = WHITE
        break
    return slide


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, color=INDIGO, align=PP_ALIGN.LEFT,
             font="Helvetica Neue", anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
             italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_filled_rect(slide, left, top, width, height, *,
                    color=PANEL_LIGHT, line_color=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = line_w
    sh.shadow.inherit = False
    return sh


def add_rounded_panel(slide, left, top, width, height, *,
                      color=PANEL_LIGHT, line_color=ROYAL, line_w=Pt(0.75),
                      adj=0.05):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                left, top, width, height)
    sh.adjustments[0] = adj
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.color.rgb = line_color
    sh.line.width = line_w
    sh.shadow.inherit = False
    return sh


def add_dot(slide, left, top, diameter=Inches(0.16), color=ROYAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_image_or_placeholder(slide, path, left, top, width, *,
                             label="image pending", height_fallback=Inches(3.4)):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), left, top, width=width)
    else:
        add_filled_rect(slide, left, top, width, height_fallback,
                        color=PANEL_LIGHT, line_color=HAIRLINE, line_w=Pt(0.75))
        add_text(slide, label, left, top + height_fallback / 2 - Inches(0.2),
                 width, Inches(0.4), size=13, color=INDIGO_DIM,
                 align=PP_ALIGN.CENTER)


def accent_bar(slide, left=Inches(0.55), top=Inches(2.0)):
    add_filled_rect(slide, left, top, Inches(0.5), Inches(0.04), color=ROYAL)


# ========== SLIDE BUILDERS ==========

def slide_02_hook(prs):
    s = add_main_slide(prs, "The Hook")
    add_text(s, "Your AI team has",
             Inches(0.55), Inches(1.25), Inches(7.5), Inches(0.9),
             size=36, bold=True, color=INDIGO)
    add_text(s, "personalities.",
             Inches(0.55), Inches(1.95), Inches(7.5), Inches(0.9),
             size=36, bold=True, color=INDIGO_SOFT)
    add_text(s, "Now give them recipes.",
             Inches(0.55), Inches(2.65), Inches(7.5), Inches(0.9),
             size=36, bold=True, color=ROYAL)
    accent_bar(s, top=Inches(3.7))
    add_text(s, "A persona is who. A skill is how. "
                "One markdown file the agent runs the same way every time, on cue.",
             Inches(0.55), Inches(3.85), Inches(7.5), Inches(1.5),
             size=15, color=INDIGO_DIM, line_spacing=1.4)

    add_image_or_placeholder(s, HERE / "hero-02-recipe.png",
                             Inches(8.6), Inches(1.3), Inches(4.3))
    return s


def slide_03_persona_vs_skill(prs):
    s = add_main_slide(prs, "Persona vs Skill")
    add_text(s, "Persona is identity.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.85),
             size=32, bold=True, color=INDIGO)
    add_text(s, "Skill is procedure. They sit on the same desk.",
             Inches(0.55), Inches(1.95), Inches(11.5), Inches(0.5),
             size=15, color=INDIGO_DIM)
    accent_bar(s, top=Inches(2.5))

    # Two columns: WHO vs HOW
    col_w = Inches(5.85)
    col_h = Inches(4.2)
    top = Inches(2.85)

    # WHO
    lx = Inches(0.55)
    add_rounded_panel(s, lx, top, col_w, col_h,
                      color=PANEL_LIGHT, line_color=ROYAL,
                      line_w=Pt(1.2), adj=0.03)
    add_filled_rect(s, lx, top, col_w, Inches(0.55), color=ROYAL)
    add_text(s, "WHO  ·  persona.md",
             lx + Inches(0.3), top + Inches(0.08),
             col_w - Inches(0.6), Inches(0.4),
             size=14, bold=True, color=WHITE, font="Menlo",
             anchor=MSO_ANCHOR.MIDDLE)
    who_lines = [
        "Aria the Scientist",
        "How she thinks. How she pushes back.",
        "Her voice. Her first five questions.",
        "Built once. Mostly stable over time.",
    ]
    y = top + Inches(0.85)
    for line in who_lines:
        add_dot(s, lx + Inches(0.4), y + Inches(0.13),
                diameter=Inches(0.12), color=ROYAL)
        add_text(s, line, lx + Inches(0.7), y,
                 col_w - Inches(0.95), Inches(0.5),
                 size=14, color=INDIGO_SOFT, line_spacing=1.3)
        y += Inches(0.65)

    # HOW
    rx = Inches(6.95)
    add_rounded_panel(s, rx, top, col_w, col_h,
                      color=PANEL_LIGHT, line_color=GOLD,
                      line_w=Pt(1.2), adj=0.03)
    add_filled_rect(s, rx, top, col_w, Inches(0.55), color=GOLD)
    add_text(s, "HOW  ·  skill.md",
             rx + Inches(0.3), top + Inches(0.08),
             col_w - Inches(0.6), Inches(0.4),
             size=14, bold=True, color=WHITE, font="Menlo",
             anchor=MSO_ANCHOR.MIDDLE)
    how_lines = [
        "weekly-funder-scan.md",
        "When to fire. What inputs to ask for.",
        "Steps in order. What good output looks like.",
        "One per repeatable task. Build a backlog.",
    ]
    y = top + Inches(0.85)
    for line in how_lines:
        add_dot(s, rx + Inches(0.4), y + Inches(0.13),
                diameter=Inches(0.12), color=GOLD)
        add_text(s, line, rx + Inches(0.7), y,
                 col_w - Inches(0.95), Inches(0.5),
                 size=14, color=INDIGO_SOFT, line_spacing=1.3)
        y += Inches(0.65)

    add_text(s, "One persona, many skills. Skills sit beside the persona file in your workspace.",
             Inches(0.55), Inches(7.25), Inches(12.2), Inches(0.4),
             size=12, color=INDIGO_DIM, align=PP_ALIGN.CENTER, italic=True)
    return s


def slide_04_master_prompt(prs):
    s = add_main_slide(prs, "One master prompt")
    add_text(s, "Paste, run, answer five questions.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=30, bold=True, color=INDIGO)
    add_text(s, "Works on ChatGPT, Claude, Gemini, Copilot.",
             Inches(0.55), Inches(2.0), Inches(11.5), Inches(0.4),
             size=13, color=INDIGO_DIM)
    accent_bar(s, top=Inches(2.45))

    # Left: the master prompt mock as a code-style panel
    panel_left = Inches(0.55)
    panel_top = Inches(2.85)
    panel_w = Inches(6.6)
    panel_h = Inches(3.95)
    add_rounded_panel(s, panel_left, panel_top, panel_w, panel_h,
                      color=PANEL_LIGHT, line_color=HAIRLINE,
                      line_w=Pt(0.75), adj=0.03)
    add_filled_rect(s, panel_left, panel_top, panel_w, Inches(0.45),
                    color=INDIGO)
    add_text(s, "master-prompt.md",
             panel_left + Inches(0.3), panel_top + Inches(0.06),
             panel_w - Inches(0.6), Inches(0.35),
             size=12, bold=True, color=WHITE, font="Menlo",
             anchor=MSO_ANCHOR.MIDDLE)

    lines = [
        ("You are the Skill Maker.", INDIGO, True),
        ("Your job is to help me write one", INDIGO_SOFT, False),
        ("SKILL.md file: a small reusable", INDIGO_SOFT, False),
        ("procedure that one of my AI personas", INDIGO_SOFT, False),
        ("can run on demand.", INDIGO_SOFT, False),
        ("", INDIGO, False),
        ("GROUND RULES", ROYAL, True),
        ("- Five questions max. One at a time.", INDIGO_SOFT, False),
        ("- The trigger description is everything.", INDIGO_SOFT, False),
        ("- Output one fenced markdown block.", INDIGO_SOFT, False),
        ("", INDIGO, False),
        ("START NOW WITH Q1.", ROYAL, True),
    ]
    y = panel_top + Inches(0.65)
    for text, color, bold in lines:
        add_text(s, text, panel_left + Inches(0.4), y,
                 panel_w - Inches(0.7), Inches(0.32),
                 size=12, bold=bold, color=color, font="Menlo",
                 line_spacing=1.1)
        y += Inches(0.27)

    # Right: five interview questions
    rx = Inches(7.55)
    rw = Inches(5.3)
    add_text(s, "The interview", rx, Inches(2.85), rw, Inches(0.4),
             size=14, bold=True, color=ROYAL)
    qs = [
        ("Q1", "Which persona owns this skill?"),
        ("Q2", "What task should the skill do?"),
        ("Q3", "When should it fire? (the trigger)"),
        ("Q4", "What does great output look like?"),
        ("Q5", "What inputs does it need?"),
    ]
    y = Inches(3.3)
    for tag, body in qs:
        add_text(s, tag, rx, y, Inches(0.9), Inches(0.32),
                 size=12, bold=True, color=ROYAL, font="Menlo")
        add_text(s, body, rx + Inches(0.9), y, rw - Inches(0.9), Inches(0.6),
                 size=14, color=INDIGO, line_spacing=1.3)
        y += Inches(0.65)
    return s


def slide_05_skill_file(prs):
    s = add_main_slide(prs, "What you get back")
    add_text(s, "One SKILL.md. Concrete, not generic.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=28, bold=True, color=INDIGO)
    accent_bar(s, top=Inches(2.0))

    # Left: a mock SKILL.md file in monospace
    panel_left = Inches(0.55)
    panel_top = Inches(2.3)
    panel_w = Inches(6.4)
    panel_h = Inches(4.55)
    add_rounded_panel(s, panel_left, panel_top, panel_w, panel_h,
                      color=PANEL_LIGHT, line_color=GOLD,
                      line_w=Pt(1.2), adj=0.03)
    add_filled_rect(s, panel_left, panel_top, panel_w, Inches(0.45),
                    color=GOLD)
    add_text(s, "weekly-funder-scan.md  ·  owner: Aria",
             panel_left + Inches(0.3), panel_top + Inches(0.06),
             panel_w - Inches(0.6), Inches(0.35),
             size=12, bold=True, color=WHITE, font="Menlo",
             anchor=MSO_ANCHOR.MIDDLE)

    body = [
        ("---", INDIGO_DIM, False),
        ("name: weekly-funder-scan", INDIGO, True),
        ("description: Use when scanning funders", INDIGO_SOFT, False),
        ("  aligned to my cause. Triggers on:", INDIGO_SOFT, False),
        ("  scan funders, any new grants?,", INDIGO_SOFT, False),
        ("  Monday research.", INDIGO_SOFT, False),
        ("owner: Aria the Scientist", INDIGO_SOFT, False),
        ("---", INDIGO_DIM, False),
        ("", INDIGO, False),
        ("# Weekly funder scan", INDIGO, True),
        ("", INDIGO, False),
        ("## When to use", ROYAL, True),
        ("- Monday morning research blocks.", INDIGO_SOFT, False),
        ("- Just before a grant deadline.", INDIGO_SOFT, False),
        ("", INDIGO, False),
        ("## Steps", ROYAL, True),
        ("1. Pull this week's open calls.", INDIGO_SOFT, False),
        ("2. Filter to my cause area.", INDIGO_SOFT, False),
        ("3. Rank by fit and deadline.", INDIGO_SOFT, False),
    ]
    y = panel_top + Inches(0.6)
    for text, color, bold in body:
        add_text(s, text, panel_left + Inches(0.35), y,
                 panel_w - Inches(0.6), Inches(0.3),
                 size=11, bold=bold, color=color, font="Menlo",
                 line_spacing=1.1)
        y += Inches(0.21)

    # Right: file format anatomy
    rx = Inches(7.45)
    rw = Inches(5.4)
    add_text(s, "Every skill file has:",
             rx, Inches(2.3), rw, Inches(0.4),
             size=14, bold=True, color=ROYAL)
    sections = [
        ("Trigger description",
         "The line that tells the AI when to use this. The single most important field."),
        ("When to use",
         "Three to five concrete situations. Plus a 'do not use when' if there's a false-trigger to avoid."),
        ("Inputs needed",
         "What the skill needs from you. Required vs optional. What to ask for if missing."),
        ("Steps",
         "The procedure. Numbered, short, no more than seven."),
        ("Output contract",
         "Length, format, what's in it, what isn't."),
        ("One worked example",
         "Real inputs in, real output out. Teaches the skill more than the steps do."),
    ]
    y = Inches(2.75)
    for h, body_t in sections:
        add_dot(s, rx, y + Inches(0.12), diameter=Inches(0.12), color=ROYAL)
        add_text(s, h, rx + Inches(0.3), y, rw - Inches(0.3), Inches(0.32),
                 size=13, bold=True, color=INDIGO)
        add_text(s, body_t, rx + Inches(0.3), y + Inches(0.32),
                 rw - Inches(0.3), Inches(0.55),
                 size=11, color=INDIGO_SOFT, line_spacing=1.25)
        y += Inches(0.66)
    return s


def slide_06_triggers(prs):
    s = add_main_slide(prs, "Triggers are everything")
    add_text(s, "A skill that never fires is worthless.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=30, bold=True, color=INDIGO)
    add_text(s, "Write triggers in your real words. The phrases you would actually say.",
             Inches(0.55), Inches(2.0), Inches(11.5), Inches(0.4),
             size=14, color=INDIGO_DIM)
    accent_bar(s, top=Inches(2.5))

    # Two stacked example panels: bad trigger / good trigger
    panel_w = Inches(12.2)
    panel_h = Inches(1.85)
    top = Inches(2.85)

    # Bad
    add_rounded_panel(s, Inches(0.55), top, panel_w, panel_h,
                      color=PANEL_LIGHT, line_color=RED_BAD,
                      line_w=Pt(1.2), adj=0.04)
    add_text(s, "TOO GENERIC",
             Inches(0.85), top + Inches(0.18),
             Inches(2.8), Inches(0.32),
             size=11, bold=True, color=RED_BAD, font="Menlo")
    add_text(s, '"a skill for funders"',
             Inches(0.85), top + Inches(0.5),
             panel_w - Inches(0.6), Inches(0.5),
             size=22, bold=True, color=INDIGO, font="Menlo")
    add_text(s, "Will never fire on its own. The model has nothing concrete to match.",
             Inches(0.85), top + Inches(1.15),
             panel_w - Inches(0.6), Inches(0.5),
             size=13, color=INDIGO_SOFT, italic=True)

    # Good
    top2 = top + panel_h + Inches(0.3)
    add_rounded_panel(s, Inches(0.55), top2, panel_w, panel_h,
                      color=PANEL_LIGHT, line_color=GREEN_OK,
                      line_w=Pt(1.2), adj=0.04)
    add_text(s, "READY TO FIRE",
             Inches(0.85), top2 + Inches(0.18),
             Inches(2.8), Inches(0.32),
             size=11, bold=True, color=GREEN_OK, font="Menlo")
    add_text(s, '"scan funders" · "any new grants?" · "Monday research"',
             Inches(0.85), top2 + Inches(0.5),
             panel_w - Inches(0.6), Inches(0.5),
             size=18, bold=True, color=INDIGO, font="Menlo")
    add_text(s, "Three real phrases. The model recognises one and runs the skill.",
             Inches(0.85), top2 + Inches(1.15),
             panel_w - Inches(0.6), Inches(0.5),
             size=13, color=INDIGO_SOFT, italic=True)

    return s


def slide_07_test(prs):
    s = add_main_slide(prs, "Test the trigger")
    add_text(s, "Save it. Load it. Say the words.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=30, bold=True, color=INDIGO)
    accent_bar(s, top=Inches(2.0))

    cards = [
        ("01", "Save",
         "Save the SKILL.md next to the persona file.\nyour-ai-workspace/skills/weekly-funder-scan.md"),
        ("02", "Load",
         "New chat. Paste the persona file plus the skill file. Or attach both."),
        ("03", "Trigger",
         "Say one of your trigger phrases verbatim.\n\"Aria, any new grants?\""),
        ("04", "Iterate",
         "Watch it run. Edit the steps if the output drifts. Skills are living documents."),
    ]
    card_w = Inches(3.0)
    card_h = Inches(3.4)
    gap = Inches(0.18)
    start_x = Inches(0.55)
    top = Inches(2.45)
    for i, (num, kind, body) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rounded_panel(s, x, top, card_w, card_h,
                          color=PANEL_LIGHT, line_color=ROYAL,
                          line_w=Pt(1.2), adj=0.04)
        add_filled_rect(s, x, top, card_w, Inches(0.5), color=ROYAL)
        add_text(s, kind, x, top + Inches(0.05), card_w, Inches(0.4),
                 size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, num, x + Inches(0.3), top + Inches(0.7),
                 Inches(1.0), Inches(0.4),
                 size=11, bold=True, color=ROYAL, font="Menlo")
        add_text(s, body, x + Inches(0.3), top + Inches(1.15),
                 card_w - Inches(0.6), card_h - Inches(1.4),
                 size=12, color=INDIGO_SOFT, line_spacing=1.35)

    add_text(s, "If it does not fire on the first try, your trigger phrases need to be closer to your real words. Edit and re-test.",
             Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.4),
             size=13, color=INDIGO_DIM, align=PP_ALIGN.CENTER, italic=True)
    return s


def slide_08_what_next(prs):
    s = add_main_slide(prs, "What you just built")
    add_text(s, "A persona who knows one task by heart.",
             Inches(0.55), Inches(1.2), Inches(7.5), Inches(0.9),
             size=28, bold=True, color=INDIGO)
    add_text(s, "One skill is a foothold. Build a backlog over time.",
             Inches(0.55), Inches(2.0), Inches(7.5), Inches(0.5),
             size=15, color=INDIGO_DIM)
    accent_bar(s, top=Inches(2.6))

    # Roadmap on the left
    items = [
        ("Today", "One persona, one skill, tested.", ROYAL),
        ("This week", "Three more skills. The tasks you do most.", ROYAL_SOFT),
        ("Next month", "A small library of skills per persona.", GOLD),
        ("Beyond", "Chains of skills become workflows.", INDIGO_SOFT),
    ]
    y = Inches(3.0)
    for tag, body, c in items:
        add_dot(s, Inches(0.55), y + Inches(0.18),
                diameter=Inches(0.16), color=c)
        add_text(s, tag, Inches(0.95), y, Inches(2.0), Inches(0.4),
                 size=14, bold=True, color=c)
        add_text(s, body, Inches(2.95), y, Inches(5.5), Inches(0.5),
                 size=14, color=INDIGO_SOFT, line_spacing=1.3)
        y += Inches(0.7)

    # Right: closing image
    add_image_or_placeholder(s, HERE / "hero-05-library.png",
                             Inches(8.6), Inches(2.5), Inches(4.3))

    return s


# ========== MAIN ==========

def remove_slide(prs, slide_index):
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    rId = sldIds[slide_index].rId
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldIds[slide_index])


def update_cover_welcome(slide, new_text):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if full.strip() == "Welcome!":
                tf = shape.text_frame
                p0 = tf.paragraphs[0]
                if p0.runs:
                    p0.runs[0].text = new_text
                    for extra_run in list(p0.runs[1:]):
                        extra_run._r.getparent().remove(extra_run._r)
                for extra_p in list(tf.paragraphs[1:]):
                    extra_p._p.getparent().remove(extra_p._p)
                return True
    return False


def build():
    prs = Presentation(str(COVER_TEMPLATE))

    while len(prs.slides) > 1:
        remove_slide(prs, len(prs.slides) - 1)

    update_cover_welcome(prs.slides[0], "Your First Skill")

    slide_02_hook(prs)
    slide_03_persona_vs_skill(prs)
    slide_04_master_prompt(prs)
    slide_05_skill_file(prs)
    slide_06_triggers(prs)
    slide_07_test(prs)
    slide_08_what_next(prs)

    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
