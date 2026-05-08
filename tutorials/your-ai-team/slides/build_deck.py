#!/usr/bin/env python3
"""
Build the Your AI Team slide deck (Lesson 8).

Modelled on tutorials/your-ai-workspace/slides/build_deck.py. 16:9, locked
Indigo Royal palette. Hero images embedded if present, labelled placeholders
otherwise. Designed to be talked over for a 12 to 15 minute video lesson.

Run: python3 build_deck.py
Output: ../your-ai-team.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
OUT = HERE.parent / "your-ai-team.pptx"
COVER_TEMPLATE = HERE.parent.parent / "cover" / "cover.pptx"

# --- Locked Indigo Royal palette ---
BG_DEEP = RGBColor(0x06, 0x0a, 0x2e)
BG_BASE = RGBColor(0x0a, 0x10, 0x45)
PANEL = RGBColor(0x12, 0x1a, 0x67)
TEXT = RGBColor(0xec, 0xf0, 0xfb)
TEXT_SOFT = RGBColor(0xc8, 0xce, 0xe4)
TEXT_DIM = RGBColor(0xa4, 0xad, 0xc8)
ACCENT = RGBColor(0x56, 0x68, 0xe0)
ACCENT_POP = RGBColor(0x00, 0x06, 0xaa)
GOLD = RGBColor(0xf4, 0xc5, 0x6b)
CYAN = RGBColor(0x5e, 0xc8, 0xe5)

# Light-mode tokens
INDIGO = RGBColor(0x1e, 0x1d, 0x4c)
INDIGO_SOFT = RGBColor(0x2b, 0x2e, 0x45)
INDIGO_DIM = RGBColor(0x5b, 0x62, 0x78)
ROYAL = RGBColor(0x00, 0x33, 0x99)
ROYAL_SOFT = RGBColor(0x6f, 0x82, 0xd4)
PANEL_LIGHT = RGBColor(0xf6, 0xf7, 0xfb)
HAIRLINE = RGBColor(0xd8, 0xdc, 0xea)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Archetype signature colours (match the headshot-prompt suggestions)
SCI_COLOR = RGBColor(0x1f, 0x4e, 0xa1)   # deep blue — Scientist
DES_COLOR = RGBColor(0xd9, 0x6b, 0x27)   # warm orange — Designer
MAK_COLOR = RGBColor(0x2f, 0x6b, 0x3d)   # forest green — Maker
COM_COLOR = RGBColor(0x8a, 0x2a, 0x3f)   # burgundy — Communicator
MAN_COLOR = RGBColor(0x44, 0x4c, 0x5a)   # charcoal grey — Manager

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _layout_by_name(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[-1]


def add_main_slide(prs, title_text):
    layout = _layout_by_name(prs, "main")
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name = "Helvetica Neue"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = WHITE
        break
    return slide


def add_blank_slide(prs):
    layout = prs.slide_layouts[-1]
    for l in prs.slide_layouts:
        if len(l.placeholders) == 0:
            layout = l
            break
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_BASE
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font="Helvetica Neue", anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_text_runs(slide, runs, left, top, width, height, *,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, props in runs:
        r = p.add_run()
        r.text = text
        r.font.name = props.get("font", "Helvetica Neue")
        r.font.size = Pt(props.get("size", 18))
        r.font.bold = props.get("bold", False)
        r.font.italic = props.get("italic", False)
        r.font.color.rgb = props.get("color", TEXT)
    return tb


def add_filled_rect(slide, left, top, width, height, *,
                    color=PANEL, line_color=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = line_w
    sh.shadow.inherit = False
    return sh


def add_rounded_panel(slide, left, top, width, height, *,
                      color=PANEL, line_color=ACCENT, line_w=Pt(0.75),
                      adj=0.05):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                left, top, width, height)
    sh.adjustments[0] = adj
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.color.rgb = line_color
    sh.line.width = line_w
    sh.shadow.inherit = False
    return sh


def add_kicker(slide, text, left, top):
    return add_text(slide, text.upper(), left, top, Inches(8), Inches(0.4),
                    size=12, bold=True, color=ACCENT, font="Helvetica Neue",
                    align=PP_ALIGN.LEFT)


def add_dot_glow(slide, left, top, diameter=Inches(0.18), color=ACCENT_POP):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_image_or_placeholder(slide, path, left, top, width, height, *,
                             label="image pending"):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), left, top, width=width, height=height)
    else:
        add_filled_rect(slide, left, top, width, height,
                        color=PANEL, line_color=ACCENT, line_w=Pt(0.75))
        add_text(slide, label, left, top + height/2 - Inches(0.2),
                 width, Inches(0.4), size=14, color=TEXT_DIM,
                 align=PP_ALIGN.CENTER)


def set_shape_alpha_fill(shape, opacity=0.6):
    from pptx.oxml.ns import qn
    sppr = shape.fill._xPr
    solidFill = sppr.find(qn('a:solidFill'))
    if solidFill is None:
        return
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        return
    alpha_val = int(opacity * 100000)
    from lxml import etree
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(alpha_val))


# ========== SLIDE BUILDERS ==========

def slide_02_hook(prs, total):
    s = add_main_slide(prs, "The Hook")
    add_text(s, "Most people use AI",
             Inches(0.55), Inches(1.25), Inches(8.0), Inches(0.9),
             size=36, bold=True, color=INDIGO, font="Helvetica Neue")
    add_text_runs(s, [
        ("like a ",
         {"size": 36, "bold": True, "color": INDIGO_SOFT}),
        ("vending machine.",
         {"size": 36, "bold": True, "color": INDIGO_DIM}),
    ], Inches(0.55), Inches(1.9), Inches(8.0), Inches(0.7))
    add_text_runs(s, [
        ("Today you build a ",
         {"size": 36, "bold": True, "color": INDIGO_SOFT}),
        ("team.",
         {"size": 36, "bold": True, "color": ROYAL}),
    ], Inches(0.55), Inches(2.5), Inches(8.0), Inches(0.7))

    add_filled_rect(s, Inches(0.55), Inches(3.4), Inches(0.5), Inches(0.05),
                    color=ROYAL)
    add_text(s, "Five personas, tailored to your life and work, "
                "built from one master prompt in twelve minutes.",
             Inches(0.55), Inches(3.55), Inches(8.0), Inches(1.2),
             size=15, color=INDIGO_DIM, font="Helvetica Neue",
             line_spacing=1.4)

    img = HERE / "hero-02-vending-vs-team.png"
    img_w = Inches(4.3)
    if img.exists():
        s.shapes.add_picture(str(img), Inches(8.6), Inches(1.3),
                             width=img_w)
    else:
        add_filled_rect(s, Inches(8.6), Inches(1.3), img_w, Inches(2.87),
                        color=PANEL_LIGHT, line_color=HAIRLINE)
        add_text(s, "vending machine → team",
                 Inches(8.6), Inches(1.3) + Inches(1.3),
                 img_w, Inches(0.4), size=12, color=INDIGO_DIM,
                 align=PP_ALIGN.CENTER)
    return s


def slide_03_framework(prs, total):
    s = add_main_slide(prs, "The Five Innovators Framework")
    add_text(s, "Every team has all five.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=34, bold=True, color=INDIGO, font="Helvetica Neue")
    add_text(s, "Each one fills a different role in how work moves.",
             Inches(0.55), Inches(2.0), Inches(11.5), Inches(0.4),
             size=13, color=INDIGO_DIM, font="Helvetica Neue")
    add_filled_rect(s, Inches(0.55), Inches(2.45), Inches(0.5), Inches(0.04),
                    color=ROYAL)

    cards = [
        ("Scientist", "Understands reality.",
         "Researches, validates, finds root causes, separates signal from noise.",
         SCI_COLOR),
        ("Designer", "Imagines what could be.",
         "Reframes problems, sketches options, gives shape and direction.",
         DES_COLOR),
        ("Maker", "Builds.",
         "Turns plans into things that exist. Iterates, ships, fixes.",
         MAK_COLOR),
        ("Communicator", "Carries the message.",
         "Writes, persuades, aligns people, creates momentum.",
         COM_COLOR),
        ("Manager", "Sustains.",
         "Plans, allocates, manages risk, keeps the system healthy.",
         MAN_COLOR),
    ]

    card_w = Inches(2.42)
    card_h = Inches(3.7)
    gap = Inches(0.13)
    start_x = Inches(0.55)
    top = Inches(2.85)
    for i, (name, tag, body, c) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rounded_panel(s, x, top, card_w, card_h,
                          color=PANEL_LIGHT, line_color=c,
                          line_w=Pt(1.5), adj=0.04)
        add_filled_rect(s, x, top, card_w, Inches(0.55), color=c)
        add_text(s, name, x, top + Inches(0.07), card_w, Inches(0.4),
                 size=15, bold=True, color=WHITE, font="Helvetica Neue",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, f"0{i+1}", x + Inches(0.2), top + Inches(0.7),
                 Inches(1.0), Inches(0.4),
                 size=11, bold=True, color=c, font="Menlo")
        add_text(s, tag, x + Inches(0.2), top + Inches(1.05),
                 card_w - Inches(0.4), Inches(0.55),
                 size=15, bold=True, color=INDIGO, font="Helvetica Neue",
                 line_spacing=1.2)
        add_text(s, body, x + Inches(0.2), top + Inches(1.85),
                 card_w - Inches(0.4), card_h - Inches(2.0),
                 size=11, color=INDIGO_SOFT, font="Helvetica Neue",
                 line_spacing=1.35)
    return s


def slide_04_master_prompt(prs, total):
    s = add_main_slide(prs, "One master prompt")
    add_text(s, "Paste, run, answer three questions.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=30, bold=True, color=INDIGO, font="Helvetica Neue")
    add_text(s, "Works on ChatGPT, Claude, Gemini, Copilot.",
             Inches(0.55), Inches(2.0), Inches(11.5), Inches(0.4),
             size=13, color=INDIGO_DIM, font="Helvetica Neue")
    add_filled_rect(s, Inches(0.55), Inches(2.45), Inches(0.5), Inches(0.04),
                    color=ROYAL)

    # Left: the master prompt mock as a code-style panel
    panel_left = Inches(0.55)
    panel_top = Inches(2.85)
    panel_w = Inches(6.6)
    panel_h = Inches(3.95)
    add_rounded_panel(s, panel_left, panel_top, panel_w, panel_h,
                      color=PANEL_LIGHT, line_color=HAIRLINE,
                      line_w=Pt(0.75), adj=0.03)
    add_filled_rect(s, panel_left, panel_top, panel_w, Inches(0.45),
                    color=INDIGO)
    add_text(s, "master-prompt.md",
             panel_left + Inches(0.3), panel_top + Inches(0.06),
             panel_w - Inches(0.6), Inches(0.35),
             size=12, bold=True, color=WHITE, font="Menlo",
             anchor=MSO_ANCHOR.MIDDLE)

    lines = [
        ("You are the AI Persona Maker.", INDIGO, True),
        ("Your job is to build my private team", INDIGO_SOFT, False),
        ("of five AI personas based on the", INDIGO_SOFT, False),
        ("Five Innovators Framework.", INDIGO_SOFT, False),
        ("", INDIGO, False),
        ("GROUND RULES", ROYAL, True),
        ("- Ask the fewest questions possible.", INDIGO_SOFT, False),
        ("- Treat me as the expert on my own life.", INDIGO_SOFT, False),
        ("- Output is markdown, five files.", INDIGO_SOFT, False),
        ("- No em dashes.", INDIGO_SOFT, False),
        ("", INDIGO, False),
        ("START NOW WITH Q1.", ROYAL, True),
    ]
    y = panel_top + Inches(0.65)
    for text, color, bold in lines:
        add_text(s, text, panel_left + Inches(0.4), y,
                 panel_w - Inches(0.7), Inches(0.32),
                 size=12, bold=bold, color=color, font="Menlo",
                 line_spacing=1.1)
        y += Inches(0.27)

    # Right: three interview questions
    rx = Inches(7.55)
    rw = Inches(5.3)
    add_text(s, "The interview", rx, Inches(2.85), rw, Inches(0.4),
             size=14, bold=True, color=ROYAL, font="Helvetica Neue")
    qs = [
        ("Q1", "Who are you and what fills your working days?"),
        ("Q2", "What does a good week look like right now?"),
        ("Q3", "How do you like to be spoken to?"),
        ("Q4 (optional)", "Should I name your team, or do you have names?"),
    ]
    y = Inches(3.3)
    for tag, body in qs:
        add_text(s, tag, rx, y, Inches(1.4), Inches(0.32),
                 size=12, bold=True, color=ROYAL, font="Menlo")
        add_text(s, body, rx + Inches(1.4), y, rw - Inches(1.4), Inches(0.7),
                 size=14, color=INDIGO, font="Helvetica Neue",
                 line_spacing=1.3)
        y += Inches(0.78)
    return s


def slide_05_persona_file(prs, total):
    s = add_main_slide(prs, "What you get back")
    add_text(s, "Five persona files. Your voice, not a template.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=28, bold=True, color=INDIGO, font="Helvetica Neue")
    add_filled_rect(s, Inches(0.55), Inches(2.0), Inches(0.5), Inches(0.04),
                    color=ROYAL)

    # Left: a single mock persona file in monospace
    panel_left = Inches(0.55)
    panel_top = Inches(2.3)
    panel_w = Inches(6.4)
    panel_h = Inches(4.5)
    add_rounded_panel(s, panel_left, panel_top, panel_w, panel_h,
                      color=PANEL_LIGHT, line_color=DES_COLOR,
                      line_w=Pt(1.5), adj=0.03)
    add_filled_rect(s, panel_left, panel_top, panel_w, Inches(0.45),
                    color=DES_COLOR)
    add_text(s, "kai.md  ·  Designer",
             panel_left + Inches(0.3), panel_top + Inches(0.06),
             panel_w - Inches(0.6), Inches(0.35),
             size=12, bold=True, color=WHITE, font="Menlo",
             anchor=MSO_ANCHOR.MIDDLE)

    body = [
        ("---", INDIGO_DIM, False),
        ("name: Kai Okonkwo", INDIGO, True),
        ("archetype: Designer", INDIGO_SOFT, False),
        ("owner: Heather", INDIGO_SOFT, False),
        ("created: 2026-05-08", INDIGO_SOFT, False),
        ("---", INDIGO_DIM, False),
        ("", INDIGO, False),
        ("# Kai, design lead for Heather", INDIGO, True),
        ("", INDIGO, False),
        ("## Who they are", ROYAL, True),
        ("Quiet, observant, sketches before speaking.", INDIGO_SOFT, False),
        ("Spent ten years redesigning patient", INDIGO_SOFT, False),
        ("intake forms. Believes a good form is a", INDIGO_SOFT, False),
        ("kind of kindness.", INDIGO_SOFT, False),
        ("", INDIGO, False),
        ("## How they push back on you", ROYAL, True),
        ('- "Is this the smallest version that', INDIGO_SOFT, False),
        ('   still moves a donor?"', INDIGO_SOFT, False),
    ]
    y = panel_top + Inches(0.6)
    for text, color, bold in body:
        add_text(s, text, panel_left + Inches(0.35), y,
                 panel_w - Inches(0.6), Inches(0.3),
                 size=11, bold=bold, color=color, font="Menlo",
                 line_spacing=1.1)
        y += Inches(0.22)

    # Right: file format anatomy
    rx = Inches(7.45)
    rw = Inches(5.4)
    add_text(s, "Every persona file has:",
             rx, Inches(2.3), rw, Inches(0.4),
             size=14, bold=True, color=ROYAL, font="Helvetica Neue")
    sections = [
        ("Who they are",
         "A real person sketch: background, temperament, why this seat."),
        ("How they think",
         "Their default moves, specific to your work."),
        ("How they push back on you",
         "Honest, useful pushback tied to your goals."),
        ("Voice",
         "How they speak, plus a sample opening line."),
        ("First five minutes",
         "Three to five questions they would put to you first."),
        ("Best used for",
         "Five concrete situations to call them in."),
        ("Headshot prompt",
         "An image-gen prompt to make their portrait."),
    ]
    y = Inches(2.75)
    for h, body_t in sections:
        add_dot_glow(s, rx, y + Inches(0.12),
                     diameter=Inches(0.12), color=ROYAL)
        add_text(s, h, rx + Inches(0.3), y, rw - Inches(0.3), Inches(0.32),
                 size=13, bold=True, color=INDIGO, font="Helvetica Neue")
        add_text(s, body_t, rx + Inches(0.3), y + Inches(0.32),
                 rw - Inches(0.3), Inches(0.5),
                 size=11, color=INDIGO_SOFT, font="Helvetica Neue",
                 line_spacing=1.25)
        y += Inches(0.59)
    return s


def slide_06_headshots(prs, total):
    s = add_main_slide(prs, "Five faces, one team")
    add_text(s, "Headshots so they feel real.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=30, bold=True, color=INDIGO, font="Helvetica Neue")
    add_text(s, "Each persona file ships with its own image-gen prompt. "
                "Generate once, save next to the .md.",
             Inches(0.55), Inches(2.0), Inches(11.5), Inches(0.4),
             size=13, color=INDIGO_DIM, font="Helvetica Neue")
    add_filled_rect(s, Inches(0.55), Inches(2.45), Inches(0.5), Inches(0.04),
                    color=ROYAL)

    # Five circular avatar slots in a row
    avatars = [
        ("Anika",  "Scientist",    SCI_COLOR),
        ("Kai",    "Designer",     DES_COLOR),
        ("Diego",  "Maker",        MAK_COLOR),
        ("Yara",   "Communicator", COM_COLOR),
        ("Rohan",  "Manager",      MAN_COLOR),
    ]
    diam = Inches(1.85)
    gap = Inches(0.45)
    total_w = 5 * diam + 4 * gap
    start_x = (SLIDE_W - total_w) / 2
    top = Inches(3.0)
    for i, (name, archetype, c) in enumerate(avatars):
        x = start_x + i * (diam + gap)
        # Halo ring
        ring = s.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.08),
                                  top - Inches(0.08),
                                  diam + Inches(0.16),
                                  diam + Inches(0.16))
        ring.fill.background()
        ring.line.color.rgb = c
        ring.line.width = Pt(2.0)
        ring.shadow.inherit = False
        # Avatar circle
        av = s.shapes.add_shape(MSO_SHAPE.OVAL, x, top, diam, diam)
        av.fill.solid()
        av.fill.fore_color.rgb = PANEL_LIGHT
        av.line.color.rgb = HAIRLINE
        av.line.width = Pt(0.75)
        av.shadow.inherit = False
        # Initial inside
        add_text(s, name[0], x, top, diam, diam,
                 size=58, bold=True, color=c, font="Helvetica Neue",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Name + archetype below
        add_text(s, name, x - Inches(0.4), top + diam + Inches(0.18),
                 diam + Inches(0.8), Inches(0.35),
                 size=15, bold=True, color=INDIGO, font="Helvetica Neue",
                 align=PP_ALIGN.CENTER)
        add_text(s, archetype, x - Inches(0.4),
                 top + diam + Inches(0.55),
                 diam + Inches(0.8), Inches(0.32),
                 size=11, color=c, font="Helvetica Neue",
                 align=PP_ALIGN.CENTER)

    # Footer note
    add_text(s, "Tip: assign each persona a signature colour so the five read "
                "as a varied team, not five versions of one face.",
             Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4),
             size=12, color=INDIGO_DIM, font="Helvetica Neue",
             align=PP_ALIGN.CENTER)
    return s


def slide_07_deploy(prs, total):
    s = add_main_slide(prs, "Deploy your team")
    add_text(s, "From markdown to teammate in 60 seconds.",
             Inches(0.55), Inches(1.2), Inches(11.5), Inches(0.9),
             size=28, bold=True, color=INDIGO, font="Helvetica Neue")
    add_filled_rect(s, Inches(0.55), Inches(2.0), Inches(0.5), Inches(0.04),
                    color=ROYAL)

    cards = [
        ("ChatGPT", "Custom GPT",
         "New GPT → paste persona as instructions → name and avatar.",
         ROYAL),
        ("Claude", "Project",
         "New project → paste persona into the project knowledge.",
         ROYAL_SOFT),
        ("Gemini", "Gem",
         "New Gem → paste persona as system instructions.",
         CYAN),
        ("Anywhere", "System prompt",
         "Open any chat, paste the file, call them by name.",
         INDIGO_SOFT),
    ]
    card_w = Inches(3.0)
    card_h = Inches(3.4)
    gap = Inches(0.18)
    start_x = Inches(0.55)
    top = Inches(2.3)
    for i, (platform, kind, body, c) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rounded_panel(s, x, top, card_w, card_h,
                          color=PANEL_LIGHT, line_color=c,
                          line_w=Pt(1.2), adj=0.04)
        add_filled_rect(s, x, top, card_w, Inches(0.5), color=c)
        add_text(s, platform, x, top + Inches(0.05), card_w, Inches(0.4),
                 size=15, bold=True, color=WHITE, font="Helvetica Neue",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, kind, x + Inches(0.3), top + Inches(0.7),
                 card_w - Inches(0.6), Inches(0.4),
                 size=18, bold=True, color=INDIGO, font="Helvetica Neue")
        add_text(s, body, x + Inches(0.3), top + Inches(1.25),
                 card_w - Inches(0.6), card_h - Inches(1.4),
                 size=12, color=INDIGO_SOFT, font="Helvetica Neue",
                 line_spacing=1.35)
        y_step = top + card_h - Inches(0.5)
        add_text(s, f"0{i+1}", x + Inches(0.3), y_step,
                 card_w - Inches(0.6), Inches(0.4),
                 size=11, bold=True, color=c, font="Menlo")

    add_text(s, "Same five .md files. Any platform. Call them by name.",
             Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.4),
             size=14, color=INDIGO_DIM, font="Helvetica Neue",
             align=PP_ALIGN.CENTER)
    return s


def slide_08_close(prs, total):
    s = add_main_slide(prs, "Closing")
    add_text(s, "You started with one assistant.",
             Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.0),
             size=40, bold=True, color=INDIGO, font="Helvetica Neue")
    add_text_runs(s, [
        ("You leave with a ",
         {"size": 56, "bold": True, "color": INDIGO}),
        ("team.",
         {"size": 56, "bold": True, "color": ROYAL}),
    ], Inches(0.55), Inches(2.7), Inches(12.2), Inches(1.2))

    add_filled_rect(s, Inches(0.55), Inches(4.2), Inches(0.5), Inches(0.05),
                    color=ROYAL)
    add_text(s, "Next: Lesson 9 — Your First Skill.",
             Inches(0.55), Inches(4.4), Inches(11.5), Inches(0.5),
             size=22, color=INDIGO_SOFT, font="Helvetica Neue")
    add_text(s, "Teach one of your five a reusable recipe they can run on demand.",
             Inches(0.55), Inches(4.95), Inches(11.5), Inches(0.5),
             size=15, color=INDIGO_DIM, font="Helvetica Neue")

    # Five tiny coloured dots as a closing motif
    dot_y = Inches(6.4)
    dot_d = Inches(0.28)
    gap = Inches(0.2)
    total_w = 5 * dot_d + 4 * gap
    dx = (SLIDE_W - total_w) / 2
    for c in [SCI_COLOR, DES_COLOR, MAK_COLOR, COM_COLOR, MAN_COLOR]:
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, dx, dot_y, dot_d, dot_d)
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
        sh.line.fill.background()
        sh.shadow.inherit = False
        dx += dot_d + gap
    return s


# ========== MAIN ==========

def remove_slide(prs, slide_index):
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    rId = sldIds[slide_index].rId
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldIds[slide_index])


def update_cover_welcome(slide, new_text):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if full.strip() == "Welcome!":
                tf = shape.text_frame
                p0 = tf.paragraphs[0]
                if p0.runs:
                    p0.runs[0].text = new_text
                    for extra_run in list(p0.runs[1:]):
                        extra_run._r.getparent().remove(extra_run._r)
                for extra_p in list(tf.paragraphs[1:]):
                    extra_p._p.getparent().remove(extra_p._p)
                return True
    return False


def build():
    prs = Presentation(str(COVER_TEMPLATE))

    while len(prs.slides) > 1:
        remove_slide(prs, len(prs.slides) - 1)

    update_cover_welcome(prs.slides[0], "Your AI Team")

    total = 8
    slide_02_hook(prs, total)
    slide_03_framework(prs, total)
    slide_04_master_prompt(prs, total)
    slide_05_persona_file(prs, total)
    slide_06_headshots(prs, total)
    slide_07_deploy(prs, total)
    slide_08_close(prs, total)

    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
