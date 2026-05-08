#!/usr/bin/env python3
"""
In-place edit of your-ai-team.pptx slide 6: replace the initial-letter
avatars with the generated headshots, masked to circles. Preserves all
other slide edits.

Run: python3 insert_headshots.py
"""

from pathlib import Path
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

HERE = Path(__file__).parent
ROOT = HERE.parent
PPTX = ROOT / "your-ai-team.pptx"
IMG_DIR = ROOT / "img"
TMP_DIR = HERE / "_circle_cache"
TMP_DIR.mkdir(exist_ok=True)

# Slide-6 layout (from build_deck.py)
SLIDE_W_IN = 13.333
DIAM_IN = 1.85
GAP_IN = 0.45
TOP_IN = 3.0
N = 5
TOTAL_W_IN = N * DIAM_IN + (N - 1) * GAP_IN
START_X_IN = (SLIDE_W_IN - TOTAL_W_IN) / 2

NAMES = ["eli", "maren", "theo", "sophie", "marcus"]


def mask_to_circle(src_path: Path, out_path: Path, size: int = 1024):
    img = Image.open(src_path).convert("RGBA")
    # Center-crop to square
    w, h = img.size
    s = min(w, h)
    left = (w - s) // 2
    top = (h - s) // 2
    img = img.crop((left, top, left + s, top + s)).resize((size, size),
                                                          Image.LANCZOS)
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size, size), fill=255)
    out = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    out.paste(img, (0, 0), mask)
    out.save(out_path, "PNG")


def main():
    # 1. Build masked PNGs
    masked = {}
    for name in NAMES:
        src = IMG_DIR / f"{name}.png"
        if not src.exists():
            raise SystemExit(f"missing {src}")
        out = TMP_DIR / f"{name}-circle.png"
        mask_to_circle(src, out)
        masked[name] = out

    # 2. Open the deck and target slide index 5 (slide 6 / headshots)
    prs = Presentation(str(PPTX))
    slide = prs.slides[5]

    # 3. Insert circular pictures, slightly inset so the existing coloured
    #    ring shows around them. Inset = 4% of diameter.
    inset = 0.04 * DIAM_IN
    pic_diam_in = DIAM_IN - 2 * inset
    for i, name in enumerate(NAMES):
        x_in = START_X_IN + i * (DIAM_IN + GAP_IN) + inset
        y_in = TOP_IN + inset
        slide.shapes.add_picture(
            str(masked[name]),
            Inches(x_in), Inches(y_in),
            width=Inches(pic_diam_in), height=Inches(pic_diam_in),
        )

    prs.save(str(PPTX))
    print(f"Updated {PPTX}")


if __name__ == "__main__":
    main()
